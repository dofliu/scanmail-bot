#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
AI 語義化檔名重命名（單層、低資訊檔名、OCR→Gemini）

為什麼這樣做：
- 掃描/相片常見的低資訊檔名不利整理；先用 OCR 或文字擷取取得「可代表內容的片段」，
  再交給 Gemini 以中文產出聚焦主題的檔名，最後加上民國日期前綴，提升可讀與可檢索性。
"""

from __future__ import annotations

import argparse
import contextlib
import dataclasses
import json
import os
import re
import sys
import time
import unicodedata
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Iterable, Tuple

# ---- 低資訊檔名判斷 ----
LOW_INFO_PATTERNS = [
    # 常見相機/手機/截圖/掃描命名
    re.compile(r"^(IMG|PXL|DSC|SCAN|SCN|PHOTO|VID|VIDEO|MOV|MVI|PANA|CIMG|GOPR|GH\d{2}|DJI)[-_ ]?\d+", re.IGNORECASE),
    re.compile(r"^(IMG|PXL)[-_ ]?\d{8}[_-]?\d{4,6}.*$", re.IGNORECASE),  # 含日期戳
    re.compile(r"^IMG-\d{8}-WA\d+.*$", re.IGNORECASE),  # WhatsApp 格式
    re.compile(r"^(WhatsApp Image|WeChat|LINE|LINE_ALBUM|Telegram)[-_ ]?.*$", re.IGNORECASE),
    re.compile(r"^(Screenshot|Screen ?Shot|截圖|螢幕截圖|螢幕擷取畫面)[-_ ]?\d{4}[-_ ]?\d{2}[-_ ]?\d{2}.*$", re.IGNORECASE),
    re.compile(r"^(掃描文件|掃描|已掃描文件)[-_ ]?\d+.*$", re.IGNORECASE),
    # 純數字或符號（無語義）
    re.compile(r"^[0-9_\-]{6,}$"),
    re.compile(r"^[0-9]{8}[_-]?[0-9]{4,6}$"),
    # UUID / 32-hex
    re.compile(r"^[a-f0-9]{32}$", re.IGNORECASE),
    re.compile(r"^[a-f0-9]{8}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{12}$", re.IGNORECASE),
    # 通用低語義單字或前綴 + 隨機碼
    re.compile(r"^(new|doc|file|scan|image|document|untitled|無標題|未命名|tmp|temp)$", re.IGNORECASE),
    re.compile(r"^(invoice|receipt|bill|statement|estimate|quote|quotation|po|purchase[_-]?order)[-_ ]?([a-z0-9]{4,}|\d{4,})([-_ ][a-z0-9]{2,})*$", re.IGNORECASE),
]

# Invoice-隨機碼-流水號（例如：Invoice-QTGTVJK9-0003）視為低資訊檔名
LOW_INFO_PATTERNS.append(
    re.compile(r"^(invoice|receipt|bill|statement|quotation|quote)[-_ ]?[A-Za-z0-9]{3,}[-_ ]?\d{3,}(?:\b.*)?$", re.IGNORECASE)
)


GENERIC_WORDS = {
    "invoice","receipt","bill","statement","document","doc","image","photo","scan","scanned",
    "file","untitled","tmp","temp","report"
}


def _compile_patterns(spec: Optional[str]) -> Optional[List[re.Pattern]]:
    if not spec:
        return None
    parts = [s for s in re.split(r",|;", spec) if s.strip()]
    compiled: List[re.Pattern] = []
    for p in parts:
        try:
            compiled.append(re.compile(p, re.IGNORECASE))
        except re.error:
            # 忽略非法 regex，避免整體失敗
            pass
    return compiled or None


def looks_low_info(stem: str, *, include: Optional[List[re.Pattern]] = None, exclude: Optional[List[re.Pattern]] = None) -> bool:
    s = stem.strip()
    if len(s) < 4:
        return True
    if include and any(p.search(s) for p in include):
        return True
    if exclude and any(p.search(s) for p in exclude):
        return False
    if any(p.search(s) for p in LOW_INFO_PATTERNS):
        return True
    # 額外啟發式：若只由 1-2 個通用詞 + 編號/雜湊片段構成，也視為低資訊
    tokens = re.split(r"[^A-Za-z0-9]+", s)
    tokens = [t for t in tokens if t]
    if not tokens:
        return True
    alpha_tokens = [t.lower() for t in tokens if re.search(r"[A-Za-z]", t)]
    digitish = [t for t in tokens if re.fullmatch(r"[A-Za-z0-9]{4,}", t)]
    if len(tokens) <= 3 and any(t in GENERIC_WORDS for t in alpha_tokens) and digitish:
        return True
    return False


# ---- 安全檔名與避免覆蓋 ----
def sanitize_filename(name: str) -> str:
    # 為什麼：Windows 保留字與非法字元會導致改名失敗或不可見
    name = unicodedata.normalize("NFC", name)
    name = "".join(ch for ch in name if unicodedata.category(ch)[0] != "C")
    name = re.sub(r"[\\/:*?\"<>|]", " ", name)
    name = re.sub(r"\s+", " ", name).strip(" .")
    reserved = {"CON","PRN","AUX","NUL","COM1","COM2","COM3","COM4","COM5","COM6","COM7","COM8","COM9","LPT1","LPT2","LPT3","LPT4","LPT5","LPT6","LPT7","LPT8","LPT9"}
    if name.upper() in reserved:
        name = f"_{name}"
    if len(name) > 120:
        name = name[:120].rstrip()
    return name or "unnamed"


def dedupe_path(target: Path) -> Path:
    if not target.exists():
        return target
    stem, suffix = target.stem, target.suffix
    parent = target.parent
    for i in range(2, 1000):
        cand = parent / f"{stem} ({i}){suffix}"
        if not cand.exists():
            return cand
    ts = int(time.time())
    return parent / f"{stem}.{ts}{suffix}"


# ---- 內容擷取（OCR/文字）----
TEXT_LIKE_EXTS = {
    ".txt",".md",".csv",".tsv",".json",".xml",".yaml",".yml",
    ".ini",".toml",".cfg",".log",
}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
XLSX_EXT = {".xls",".xlsx"}
IMG_EXT = {".png",".jpg",".jpeg",".bmp",".tif",".tiff"}
PDF_EXT = {".pdf"}


def extract_text_snippet(path: Path, max_chars: int = 2000) -> Optional[str]:
    suf = path.suffix.lower()

    # 純文字：低成本直接讀
    if suf in TEXT_LIKE_EXTS:
        for enc in ("utf-8","utf-8-sig","cp950","big5","latin1"):
            with contextlib.suppress(Exception):
                txt = path.read_text(encoding=enc, errors="ignore")
                txt = re.sub(r"\s+"," ", txt).strip()
                if txt:
                    return txt[:max_chars]
        return None

    # PDF：若可抽取文字就抽；掃描型 PDF 沒內嵌文字則回 None（避免誤用）
    if suf in PDF_EXT:
        with contextlib.suppress(Exception):
            import PyPDF2  # type: ignore
            parts: List[str] = []
            with open(path, "rb") as f:
                r = PyPDF2.PdfReader(f)
                for page in r.pages[:5]:
                    t = page.extract_text() or ""
                    parts.append(t)
                    if sum(len(x) for x in parts) > max_chars:
                        break
            text = re.sub(r"\s+"," ", " ".join(parts)).strip()
            return text[:max_chars] if text else None

    # 影像：OCR（需 Tesseract 與 pytesseract）
    if suf in IMG_EXT:
        with contextlib.suppress(Exception):
            from PIL import Image  # type: ignore
            import pytesseract  # type: ignore
            img = Image.open(path)
            t = pytesseract.image_to_string(img, lang="chi_tra+eng")
            t = re.sub(r"\s+"," ", t).strip()
            return t[:max_chars] if t else None
        return None

    # DOCX
    if suf in DOCX_EXT:
        with contextlib.suppress(Exception):
            import docx  # type: ignore
            d = docx.Document(str(path))
            text = re.sub(r"\s+"," ", " ".join(p.text for p in d.paragraphs)).strip()
            return text[:max_chars] if text else None

    # PPTX
    if suf in PPTX_EXT:
        with contextlib.suppress(Exception):
            from pptx import Presentation  # type: ignore
            prs = Presentation(str(path))
            texts: List[str] = []
            for slide in prs.slides[:10]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            text = re.sub(r"\s+"," ", " ".join(texts)).strip()
            return text[:max_chars] if text else None

    # XLSX（只抓前 50 列）
    if suf in XLSX_EXT:
        with contextlib.suppress(Exception):
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            ws = wb.active
            vals: List[str] = []
            for row in ws.iter_rows(min_row=1, max_row=50, values_only=True):
                line = " ".join(str(c) for c in row if c is not None)
                if line.strip():
                    vals.append(line)
                if sum(len(x) for x in vals) > max_chars:
                    break
            text = re.sub(r"\s+"," ", " ".join(vals)).strip()
            return text[:max_chars] if text else None

    # 其他：略過，避免誤判
    return None


# ---- Gemini 介面 ----
@dataclasses.dataclass
class GeminiConfig:
    api_key: str
    model: str = os.environ.get("GEMINI_MODEL", os.environ.get("GOOGLE_MODEL", "gemini-3-flash-preview"))
    timeout: int = int(os.environ.get("GEMINI_TIMEOUT", "60"))


def gemini_generate(prompt: str, cfg: GeminiConfig) -> str:
    """
    為什麼：使用官方 REST v1beta generateContent，輸入文字提示，
    僅回傳第一候選的文字內容；限制輸出長度交由提示控制。
    """
    import urllib.request
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{cfg.model}:generateContent?key={cfg.api_key}"
    payload = {
        "contents": [
            {"parts": [{"text": prompt}]}
        ]
    }
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(url, data=data, headers={"Content-Type": "application/json"}, method="POST")
    with urllib.request.urlopen(req, timeout=cfg.timeout) as resp:
        body = resp.read().decode("utf-8", errors="ignore")
        js = json.loads(body)
    cands = js.get("candidates") or []
    if not cands:
        raise RuntimeError("Gemini 無回覆候選")
    parts = ((cands[0].get("content") or {}).get("parts") or [])
    texts = [p.get("text","") for p in parts if isinstance(p, dict)]
    out = " ".join(texts).strip()
    return out


PROMPT_TEMPLATE = (
    "你是一個檔案整理助手。根據以下內容摘要，產生一個最合適的『中文檔名主體』（不含副檔名）。"
    "規則：1) 最多 40 字；2) 專有名詞保留英文；3) 涵蓋主題與關鍵詞；"
    "4) 不加引號或特殊符號；5) 只輸出檔名本身。\n\n"
    "[內容]\n{snippet}\n\n"
    "只輸出檔名："
)


def suggest_filename_with_gemini(snippet: str, cfg: GeminiConfig) -> str:
    raw = gemini_generate(PROMPT_TEMPLATE.format(snippet=snippet[:1800]), cfg)
    raw = raw.strip().strip("'\"“”`，。！!？?：:；;[]（）(){}《》<> ")
    return sanitize_filename(raw)


# ---- 日期前綴（民國年）----
def roc_yyyMMdd_from_mtime(p: Path) -> str:
    # 為什麼用 mtime：更能代表檔案最後被確認的實際時間；可避免改用今日造成資訊汙染
    ts = p.stat().st_mtime
    dt = datetime.fromtimestamp(ts)
    roc_year = max(1, dt.year - 1911)
    return f"{roc_year:03d}{dt:%m%d}"


# ---- 主流程（單層、低資訊）----
def main(argv: Optional[List[str]] = None) -> int:
    ap = argparse.ArgumentParser(description="AI 語義化檔名重命名（本層、低資訊、OCR→Gemini）")
    ap.add_argument("--dir", default=".", help="目標資料夾（只掃描本層）")
    ap.add_argument("--rename", action="store_true", help="實際改名（預設為 dry-run 預覽）")
    ap.add_argument("--only-exts", default=None, help="僅處理指定副檔名（逗號分隔，如：pdf,docx,png）")
    ap.add_argument("--include-pattern", default=None, help="額外視為低資訊的 regex（可用逗號分隔多個）")
    ap.add_argument("--exclude-pattern", default=None, help="強制視為非低資訊的 regex（可用逗號分隔多個）")
    ap.add_argument("--google-api-key", dest="google_api_key", default=os.environ.get("GOOGLE_API_KEY") or os.environ.get("GEMINI_API_KEY"), help="Google API Key（亦可用環境變數）")
    ap.add_argument("--model", default=os.environ.get("GEMINI_MODEL", os.environ.get("GOOGLE_MODEL", "gemini-3-flash-preview")), help="Gemini 模型名稱")
    ap.add_argument("--debug", action="store_true", help="輸出每個檔案的判斷原因（協助調整規則）")
    args = ap.parse_args(argv)

    if not args.google_api_key:
        print("[錯誤] 缺少 Google API Key，請以 --google-api-key 或環境變數 GOOGLE_API_KEY / GEMINI_API_KEY 提供", file=sys.stderr)
        return 2

    root = Path(args.dir).expanduser().resolve()
    if not root.exists() or not root.is_dir():
        print(f"[錯誤] 目標資料夾不存在：{root}", file=sys.stderr)
        return 2

    only_exts = None if not args.only_exts else {"." + s.strip().lower().lstrip(".") for s in args.only_exts.split(",") if s.strip()}

    files = [p for p in root.glob("*") if p.is_file()]
    if not files:
        print("[資訊] 沒有找到檔案。")
        return 0

    cfg = GeminiConfig(api_key=args.google_api_key, model=args.model)
    incl = _compile_patterns(args.include_pattern)
    excl = _compile_patterns(args.exclude_pattern)

    print(f"[資訊] 掃描 {root}（本層），檔案數：{len(files)}，模式：{'rename' if args.rename else 'dry-run'}，模型：{cfg.model}")

    changed = skipped = failed = 0

    for src in files:
        if only_exts and src.suffix.lower() not in only_exts:
            if args.debug:
                print(f"[debug] skip-ext  {src.name} -> {src.suffix.lower()} 不在 {sorted(only_exts)}")
            continue
        if not looks_low_info(src.stem, include=incl, exclude=excl):
            if args.debug:
                print(f"[debug] not-low  {src.name} -> 檔名看起來已有語義（可用 --include-pattern 覆寫）")
            continue
        else:
            if args.debug:
                print(f"[debug] candidate {src.name} -> 低資訊檔名（規則/覆寫），嘗試擷取內容")

        snippet = extract_text_snippet(src)
        if not snippet:
            skipped += 1
            print(f"[略過] {src.name}（no-content:low-info）")
            continue

        try:
            base = suggest_filename_with_gemini(snippet, cfg)
        except Exception as e:
            failed += 1
            print(f"[失敗] {src.name} -> LLM：{e}")
            continue

        # 與原名相同則略過
        if base and base.lower() == src.stem.lower():
            skipped += 1
            print(f"[略過] {src.name}（no-change:low-info）")
            continue

        prefix = roc_yyyMMdd_from_mtime(src)
        new_stem = f"{prefix}-{base}"
        new_name = sanitize_filename(new_stem) + src.suffix
        dst = dedupe_path(src.with_name(new_name))

        if args.rename:
            try:
                os.replace(src, dst)
                changed += 1
                print(f"[改名] {src.name} -> {dst.name}")
            except Exception as e:
                failed += 1
                print(f"[失敗] {src.name} -> {dst.name}：{e}")
        else:
            changed += 1
            print(f"[預覽] {src.name} -> {dst.name}")

    print(f"[總結] 可改名：{changed}，略過：{skipped}，失敗：{failed}")
    if not args.rename and changed > 0:
        print("[提示] 以上為預覽。加入 --rename 以套用。")
    return 0 if failed == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
